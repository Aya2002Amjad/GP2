import io
import docx
import pdfplumber
from pptx import Presentation 

def process_document(file_bytes,file_extension):
    if file_extension.lower()=='.pdf':
        return process_pdf(file_bytes)
    elif file_extension.lower()=='.docx':
        return process_word(file_bytes)
    elif  file_extension.lower()=='.txt':
        return process_text(file_bytes)
    elif file_extension == '.pptx':
        return process_pptx(file_bytes)
    else:
        return "Unsupported file format"



def process_pdf(fily_bytes):
    with pdfplumber.open(io.BytesIO(fily_bytes)) as pdf:
        text=''
        for page in pdf.pages:
            text+=page.extract_text()
    return text

def process_word(file_bytes):
    doc=docx.Document(io.BytesIO(file_bytes))
    text=''
    for paragraph in doc.paragraphs:
        text+=paragraph.text+'\n'
    return text


def process_text(file_bytes):
    return file_bytes.decode('utf-8')

def process_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text